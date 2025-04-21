import fitz
from docx import Document
from pptx import Presentation
from PIL import Image
import io

def extract_text_from_pdf(file_path):
    try:
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text
    except Exception as e:
        return None, f"Error extracting text from PDF: {e}"

def extract_text_from_word(file_path):
    try:
        doc = Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except Exception as e:
        return None, f"Error extracting text from Word document: {e}"

def extract_text_from_ppt(file_path):
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
                        text += "\n"
        return text
    except Exception as e:
        return None, f"Error extracting text from PowerPoint: {e}"

def process_image_and_get_prompt(file_path, question):
    try:
        with open(file_path, 'rb') as image_file:
            image_data = image_file.read()
        image_part = {"mime_type": Image.MIME[Image.open(io.BytesIO(image_data)).format], "data": image_data}
        prompt_parts = [image_part, question]
        return prompt_parts, None
    except Exception as e:
        return None, f"Error processing image: {e}"